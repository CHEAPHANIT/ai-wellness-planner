from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "NutriAI_Class_Presentation_Designed.pptx"

INK = RGBColor(18, 24, 38)
MUTED = RGBColor(91, 104, 124)
GREEN = RGBColor(18, 166, 105)
GREEN_DARK = RGBColor(8, 122, 77)
MINT = RGBColor(229, 248, 239)
PAPER = RGBColor(248, 251, 249)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(206, 225, 215)
BLUE = RGBColor(62, 127, 248)
BLUE_SOFT = RGBColor(234, 242, 255)
AMBER = RGBColor(246, 147, 47)
AMBER_SOFT = RGBColor(255, 243, 224)
CORAL = RGBColor(239, 96, 94)
CORAL_SOFT = RGBColor(255, 236, 235)
PURPLE = RGBColor(111, 101, 230)
PURPLE_SOFT = RGBColor(241, 239, 255)
DARK_PANEL = RGBColor(8, 63, 49)


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, x, y, w, h, lines, size=14, color=INK, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        if bullet:
            p.text = f"- {line}"
        else:
            p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def fill(shape, color, line=LINE):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line


def add_round_rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, line)
    return shape


def add_slide_number(slide, number):
    add_text(slide, 11.83, 7.02, 0.55, 0.22, f"{number:02}", 8, GREEN, True, PP_ALIGN.RIGHT)


def add_brand(slide, dark=False):
    color = WHITE if dark else GREEN_DARK
    leaf = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.58), Inches(0.35), Inches(0.34), Inches(0.34))
    fill(leaf, GREEN if dark else MINT, GREEN if dark else MINT)
    add_text(slide, 0.67, 0.42, 0.16, 0.14, "N", 9, color, True, PP_ALIGN.CENTER)
    add_text(slide, 1.0, 0.38, 1.3, 0.2, "NutriAI", 13, color, True)


def add_header(slide, number, eyebrow, title):
    add_brand(slide)
    add_text(slide, 0.65, 0.95, 2.4, 0.22, eyebrow.upper(), 9, GREEN_DARK, True)
    add_text(slide, 0.65, 1.22, 10.5, 0.55, title, 25, INK, True)
    add_slide_number(slide, number)


def add_chip(slide, x, y, text, color=GREEN, soft=MINT, w=1.65):
    chip = add_round_rect(slide, x, y, w, 0.42, soft, soft)
    add_text(slide, x + 0.08, y + 0.12, w - 0.16, 0.12, text, 9, color, True, PP_ALIGN.CENTER)
    return chip


def add_card(slide, x, y, w, h, title, body, accent=GREEN, soft=MINT):
    add_round_rect(slide, x, y, w, h, WHITE, LINE)
    add_round_rect(slide, x + 0.22, y + 0.24, 0.48, 0.48, soft, soft)
    add_text(slide, x + 0.35, y + 0.36, 0.22, 0.12, title[0].upper(), 10, accent, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.86, y + 0.22, w - 1.05, 0.28, title, 13, INK, True)
    add_text(slide, x + 0.86, y + 0.58, w - 1.05, h - 0.68, body, 9.5, MUTED)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_PANEL
    add_brand(slide, dark=True)
    add_text(slide, 0.75, 1.35, 3.3, 0.25, "CLASS PRESENTATION", 11, RGBColor(139, 222, 184), True)
    add_text(slide, 0.75, 1.92, 5.2, 0.75, "NutriAI", 56, WHITE, True)
    add_text(slide, 0.78, 2.85, 6.8, 0.42, "AI Nutrition and Meal Planner", 24, RGBColor(215, 238, 226), False)
    add_multiline(
        slide,
        0.8,
        3.58,
        5.8,
        0.85,
        ["Personalized meal plans", "Allergy-aware recommendations", "Food, water, weight, grocery and AI assistant workflows"],
        12,
        RGBColor(194, 224, 209),
        True,
    )
    add_round_rect(slide, 7.35, 1.1, 4.8, 4.9, RGBColor(12, 82, 63), RGBColor(42, 120, 91))
    add_text(slide, 7.75, 1.55, 3.7, 0.32, "What the system does", 18, WHITE, True)
    for i, (label, value, color) in enumerate(
        [
            ("01", "Collect profile and goal data", GREEN),
            ("02", "Filter unsafe foods using allergies", BLUE),
            ("03", "Score meal candidates by nutrition targets", AMBER),
            ("04", "Turn plans into groceries and logs", PURPLE),
        ]
    ):
        y = 2.1 + i * 0.82
        add_round_rect(slide, 7.75, y, 0.5, 0.38, color, color)
        add_text(slide, 7.88, y + 0.11, 0.24, 0.1, label, 8, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, 8.45, y + 0.08, 3.0, 0.18, value, 11, RGBColor(231, 244, 238), True)
    add_text(slide, 0.8, 6.55, 4.6, 0.24, "Presented by: Group Members", 10, RGBColor(194, 224, 209), True)
    add_slide_number(slide, 1)


def add_members(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 2, "Team", "Group Members")
    names = ["Member 1 Name", "Member 2 Name", "Member 3 Name"]
    colors = [(GREEN, MINT), (BLUE, BLUE_SOFT), (AMBER, AMBER_SOFT)]
    for i, name in enumerate(names):
        x = 1.05 + i * 3.9
        add_round_rect(slide, x, 2.15, 3.05, 3.75, WHITE, LINE)
        accent, soft = colors[i]
        oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.78), Inches(2.58), Inches(1.5), Inches(1.5))
        fill(oval, soft, accent)
        add_text(slide, x + 0.94, 3.03, 1.18, 0.2, "PHOTO", 12, accent, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.4, 4.38, 2.25, 0.25, name, 16, INK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.45, 4.82, 2.15, 0.24, "Role / Responsibility", 10, MUTED, False, PP_ALIGN.CENTER)
        add_chip(slide, x + 0.66, 5.26, "replace image", accent, soft, 1.72)


def add_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 3, "Roadmap", "Contents / Agenda")
    items = [
        "Project overview",
        "Problem statement",
        "AI technologies and models used",
        "System architecture",
        "Key features",
        "Result and demo",
        "Challenges and lessons learned",
    ]
    for i, item in enumerate(items):
        y = 2.02 + i * 0.56
        add_round_rect(slide, 1.0, y, 0.38, 0.32, GREEN if i % 2 == 0 else BLUE, GREEN if i % 2 == 0 else BLUE)
        add_text(slide, 1.12, y + 0.09, 0.14, 0.1, str(i + 1), 8, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, 1.6, y + 0.06, 5.6, 0.18, item, 14, INK, True)
    add_round_rect(slide, 8.25, 2.08, 3.7, 3.7, DARK_PANEL, DARK_PANEL)
    add_text(slide, 8.65, 2.55, 2.9, 0.25, "Presentation flow", 18, WHITE, True)
    add_multiline(
        slide,
        8.65,
        3.08,
        2.8,
        1.35,
        ["Why the project matters", "How the system works", "What the demo proves"],
        12,
        RGBColor(213, 238, 226),
        True,
    )
    add_chip(slide, 8.68, 4.92, "full-stack + AI", GREEN, RGBColor(18, 93, 70), 2.4)


def add_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 4, "Overview", "NutriAI Is a Connected Nutrition Planning App")
    cards = [
        ("Personalized", "Uses profile, goal, activity, preference and allergy data.", GREEN, MINT),
        ("Actionable", "Generates meals that can become grocery lists and daily logs.", BLUE, BLUE_SOFT),
        ("Explainable", "Uses scoring, rules and formulas so recommendations are easier to trust.", AMBER, AMBER_SOFT),
        ("Reliable", "Optional LLM support with a local fallback when providers are unavailable.", PURPLE, PURPLE_SOFT),
    ]
    for i, card in enumerate(cards):
        row, col = divmod(i, 2)
        add_card(slide, 0.85 + col * 5.65, 2.0 + row * 1.75, 5.1, 1.25, *card)
    add_round_rect(slide, 1.5, 5.8, 10.2, 0.62, GREEN_DARK, GREEN_DARK)
    add_text(slide, 1.75, 6.0, 9.7, 0.18, "Main idea: move from personal health data to practical food decisions.", 14, WHITE, True, PP_ALIGN.CENTER)


def add_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 5, "Problem", "Healthy Meal Planning Has Too Many Moving Parts")
    problems = [
        ("Complex targets", "Calories, macros, activity, goals and eating habits must fit together."),
        ("Allergy risk", "Generic plans may recommend unsafe ingredients."),
        ("Decision fatigue", "Users know the goal but do not know what to eat today."),
        ("Disconnected tools", "Planning, groceries and tracking often live in separate apps."),
    ]
    for i, (title, body) in enumerate(problems):
        add_card(slide, 0.9 + (i % 2) * 5.75, 2.0 + (i // 2) * 1.75, 5.15, 1.28, title, body, [CORAL, AMBER, BLUE, GREEN][i], [CORAL_SOFT, AMBER_SOFT, BLUE_SOFT, MINT][i])
    add_text(slide, 1.0, 5.92, 10.9, 0.28, "NutriAI treats meal recommendation as a safety-aware, context-aware workflow instead of a random recipe list.", 15, GREEN_DARK, True, PP_ALIGN.CENTER)


def add_ai(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 6, "AI Stack", "AI Technologies / Models Used")
    items = [
        ("Meal recommender", "Heuristic combinational search and nutrition scoring", GREEN, MINT),
        ("Calorie model", "Mifflin-St Jeor BMR plus activity and goal adjustment", BLUE, BLUE_SOFT),
        ("Substitute ranking", "Nutrition-distance comparison for safe alternatives", AMBER, AMBER_SOFT),
        ("Health risk", "Transparent BMI, habits and exercise rule checks", CORAL, CORAL_SOFT),
        ("Chat assistant", "OpenAI or Ollama integration with local fallback", PURPLE, PURPLE_SOFT),
        ("Photo helper", "Conservative low-confidence prototype, avoids overclaiming", GREEN_DARK, MINT),
    ]
    for i, item in enumerate(items):
        row, col = divmod(i, 3)
        add_card(slide, 0.65 + col * 4.15, 1.95 + row * 1.82, 3.65, 1.35, *item)


def add_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 7, "Architecture", "System Architecture")
    layers = [
        ("Flutter Web", "Screens, forms, charts, navigation", GREEN, MINT),
        ("FastAPI Backend", "REST API, JWT auth, validation, AI services", BLUE, BLUE_SOFT),
        ("PostgreSQL", "Users, profiles, foods, plans, logs, allergies", AMBER, AMBER_SOFT),
        ("Deployment", "Docker locally, Vercel and Neon publicly", PURPLE, PURPLE_SOFT),
    ]
    x_positions = [0.75, 3.8, 6.85, 9.9]
    for i, (title, body, accent, soft) in enumerate(layers):
        x = x_positions[i]
        add_round_rect(slide, x, 2.45, 2.35, 2.0, WHITE, LINE)
        add_round_rect(slide, x + 0.66, 2.78, 1.0, 0.48, soft, soft)
        add_text(slide, x + 0.82, 2.93, 0.68, 0.12, str(i + 1), 11, accent, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.22, 3.45, 1.9, 0.22, title, 14, INK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.25, 3.86, 1.85, 0.4, body, 8.6, MUTED, False, PP_ALIGN.CENTER)
        if i < 3:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + 2.38),
                Inches(3.42),
                Inches(x + 3.02),
                Inches(3.42),
            )
            connector.line.color.rgb = GREEN_DARK
            connector.line.width = Pt(2)
    add_round_rect(slide, 1.45, 5.48, 10.45, 0.55, DARK_PANEL, DARK_PANEL)
    add_text(slide, 1.62, 5.65, 10.1, 0.15, "Request flow: Flutter sends JSON requests -> FastAPI runs services -> PostgreSQL stores and returns personalized results", 11, WHITE, True, PP_ALIGN.CENTER)


def add_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 8, "Features", "Key Features")
    features = [
        ("Profile", "BMI, BMR, targets, preferences"),
        ("Allergies", "Restrictions used before recommendations"),
        ("Meal planner", "Daily and weekly plan generation"),
        ("Grocery", "Estimated cost and export workflow"),
        ("Tracking", "Food, water, weight and progress"),
        ("AI assistant", "Nutrition Q&A, substitutes, risk check"),
    ]
    accents = [(GREEN, MINT), (BLUE, BLUE_SOFT), (AMBER, AMBER_SOFT), (PURPLE, PURPLE_SOFT), (CORAL, CORAL_SOFT), (GREEN_DARK, MINT)]
    for i, (title, body) in enumerate(features):
        row, col = divmod(i, 3)
        accent, soft = accents[i]
        add_card(slide, 0.65 + col * 4.15, 1.95 + row * 1.82, 3.65, 1.35, title, body, accent, soft)


def add_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 9, "Results", "Result and Demo Flow")
    steps = [
        "Register / log in",
        "Complete profile",
        "Add allergy",
        "Generate meal plan",
        "Export grocery list",
        "Ask AI assistant",
    ]
    for i, step in enumerate(steps):
        x = 0.8 + i * 2.02
        add_round_rect(slide, x, 2.3, 1.45, 1.32, [GREEN, BLUE, AMBER, PURPLE, CORAL, GREEN_DARK][i], [GREEN, BLUE, AMBER, PURPLE, CORAL, GREEN_DARK][i])
        add_text(slide, x + 0.17, 2.55, 1.1, 0.18, f"0{i + 1}", 11, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.17, 2.93, 1.1, 0.27, step, 9.5, WHITE, True, PP_ALIGN.CENTER)
    add_round_rect(slide, 1.0, 4.65, 5.1, 1.25, WHITE, LINE)
    add_text(slide, 1.3, 4.9, 4.4, 0.2, "Expected demo result", 14, INK, True)
    add_multiline(slide, 1.3, 5.25, 4.5, 0.45, ["Generated plans reflect targets", "Allergy foods are excluded", "Grocery and tracking workflows connect"], 9.5, MUTED, True)
    add_round_rect(slide, 6.9, 4.65, 4.9, 1.25, WHITE, LINE)
    add_text(slide, 7.2, 4.9, 4.2, 0.2, "Project result", 14, INK, True)
    add_text(slide, 7.2, 5.28, 4.1, 0.3, "A working full-stack application with backend service tests and Flutter widget tests included.", 10, MUTED)


def add_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, 10, "Reflection", "Challenges and Lessons Learned")
    lessons = [
        ("Allergy safety", "Filter unsafe foods before scoring candidates."),
        ("Fallback design", "External AI provider failures should not stop the app."),
        ("Honest AI", "Low-confidence photo support should say unknown instead of guessing."),
        ("Data quality", "Profile inputs must be real before showing calculated results."),
        ("Testing", "Meal planning touches many tables and needs end-to-end checks."),
        ("Explainability", "Nutrition advice should be practical, transparent and cautious."),
    ]
    for i, (title, body) in enumerate(lessons):
        row, col = divmod(i, 2)
        add_card(slide, 0.9 + col * 5.75, 1.88 + row * 1.42, 5.15, 1.05, title, body, [GREEN, BLUE, AMBER, PURPLE, CORAL, GREEN_DARK][i], [MINT, BLUE_SOFT, AMBER_SOFT, PURPLE_SOFT, CORAL_SOFT, MINT][i])


def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_PANEL
    add_brand(slide, dark=True)
    add_text(slide, 0.85, 1.35, 10.6, 0.55, "NutriAI turns personal health data into practical food decisions.", 30, WHITE, True)
    add_text(slide, 0.9, 2.25, 9.4, 0.32, "Personalize. Plan. Shop. Track. Improve.", 21, RGBColor(179, 234, 208), True)
    for i, word in enumerate(["Profile", "Allergy", "Meals", "Grocery", "Tracking", "Assistant"]):
        add_chip(slide, 0.95 + i * 1.86, 3.65, word, WHITE, RGBColor(15, 82, 63), 1.45)
    add_text(slide, 4.9, 5.72, 2.6, 0.35, "Thank You", 25, RGBColor(137, 224, 185), True, PP_ALIGN.CENTER)
    add_slide_number(slide, 11)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_members(prs)
    add_agenda(prs)
    add_overview(prs)
    add_problem(prs)
    add_ai(prs)
    add_architecture(prs)
    add_features(prs)
    add_demo(prs)
    add_challenges(prs)
    add_closing(prs)

    prs.save(OUTPUT)


def main():
    build()
    print(f"Created: {OUTPUT}")
    print("Slides: 11")


if __name__ == "__main__":
    main()
