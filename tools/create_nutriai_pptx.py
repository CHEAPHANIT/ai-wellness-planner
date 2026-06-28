from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "NutriAI_Class_Presentation.pptx"

INK = RGBColor(16, 35, 27)
MUTED = RGBColor(96, 112, 106)
GREEN = RGBColor(19, 166, 106)
GREEN_DARK = RGBColor(8, 122, 77)
MINT = RGBColor(234, 248, 241)
PAPER = RGBColor(247, 251, 249)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(61, 123, 250)
BLUE_SOFT = RGBColor(234, 241, 255)
AMBER = RGBColor(244, 166, 41)
AMBER_SOFT = RGBColor(255, 244, 223)
CORAL = RGBColor(238, 107, 98)
CORAL_SOFT = RGBColor(255, 240, 238)
PURPLE = RGBColor(120, 103, 232)
PURPLE_SOFT = RGBColor(240, 238, 255)
LINE = RGBColor(203, 228, 215)


SLIDES = [
    {
        "type": "cover",
        "title": "NutriAI",
        "subtitle": "AI Nutrition and Meal Planner",
        "items": ["Presented by: Your Name", "Course: Your Course Name", "June 2026"],
    },
    {
        "title": "A full-stack AI nutrition planning system",
        "subtitle": "Project overview",
        "items": [
            ("Full stack", "Flutter frontend and FastAPI backend"),
            ("Personalized", "Uses profile, goals, allergies and preferences"),
            ("Actionable", "Turns plans into grocery and tracking workflows"),
            ("Hybrid AI", "Combines scoring, rules, formulas and optional LLMs"),
            ("Local foods", "Includes Cambodian and familiar food options"),
            ("Docker ready", "Runs with PostgreSQL, backend and web frontend"),
        ],
    },
    {
        "title": "Healthy meal planning has too many moving parts",
        "subtitle": "Problem statement",
        "items": [
            ("Complex targets", "Users must balance calories, macros, activity and health goals."),
            ("Safety risk", "Generic plans can recommend unsafe foods for allergies."),
            ("Disconnected tools", "Planning, shopping and tracking are often separated."),
        ],
    },
    {
        "title": "NutriAI connects the complete daily nutrition workflow",
        "subtitle": "Proposed solution",
        "items": [
            ("Profile", "Collect body data, activity and goal"),
            ("Protect", "Store allergies and remove unsafe foods"),
            ("Recommend", "Generate daily or weekly meal plans"),
            ("Prepare", "Export plans to grocery lists"),
            ("Track", "Log food, water and weight"),
            ("Guide", "Answer nutrition questions with context"),
        ],
    },
    {
        "title": "A simple full-stack architecture",
        "subtitle": "System architecture",
        "items": [
            ("Flutter", "Frontend", "Screens, forms, charts"),
            ("FastAPI", "Backend", "REST API, auth, AI logic"),
            ("PostgreSQL", "Database", "Users, foods, plans, logs"),
            ("Vercel + Neon", "Deployment", "Hosted web, API and database"),
        ],
    },
    {
        "title": "NutriAI uses a hybrid set of AI techniques",
        "subtitle": "AI technologies and models",
        "items": [
            ("Meal recommendation", "Combinational search and heuristic scoring"),
            ("Calorie targets", "Mifflin-St Jeor BMR and activity formulas"),
            ("Food substitutes", "Nutrition-distance ranking"),
            ("Health screening", "Transparent rule-based risk scoring"),
            ("Chatbot", "Ollama, OpenAI or local fallback"),
            ("Photo assistance", "Conservative low-confidence prototype"),
        ],
    },
    {
        "title": "Key features support planning, shopping and tracking",
        "subtitle": "Application features",
        "items": [
            ("Authentication", "Register, login, logout and JWT sessions"),
            ("Profile", "BMI, BMR, activity and calorie targets"),
            ("Meal plans", "Daily and weekly recommendation generation"),
            ("Groceries", "Export plans with estimated costs"),
            ("Tracking", "Food, water, weight and nutrition totals"),
            ("AI assistant", "Chat, substitutes, risk check and photo support"),
        ],
    },
    {
        "title": "Meal recommendation is the core intelligent feature",
        "subtitle": "AI detail 1 - heuristic search",
        "items": [
            ("1", "Load foods", "Read available foods from PostgreSQL."),
            ("2", "Remove allergens", "Exclude foods matching user allergy records."),
            ("3", "Apply rules", "Use goal, preference, category and budget signals."),
            ("4", "Build candidates", "Create combinations containing one to three foods."),
            ("5", "Score", "Compare calories, protein, carbs, fat and category fit."),
            ("6", "Select", "Return the lowest-scoring meal candidate."),
        ],
        "footer": "score = calorie gap + macro gaps + category penalty + budget penalty",
    },
    {
        "title": "The assistant and support services stay explainable",
        "subtitle": "AI detail 2 - chatbot and supporting logic",
        "items": [
            ("Chatbot", "Uses goal, targets, allergies and today's calories."),
            ("Fallback", "Local rules respond when Ollama or OpenAI is unavailable."),
            ("Substitutes", "Rank alternatives by calories, macros and category distance."),
            ("Safety", "Risk and photo features explain limits instead of overclaiming."),
        ],
    },
    {
        "title": "The demo proves the main AI workflow",
        "subtitle": "Results and demo flow",
        "items": [
            ("1", "Register or log in", "Use the test account credentials."),
            ("2", "Complete profile", "Save body data, goal and preference."),
            ("3", "Add milk allergy", "Create a restriction used by recommendations."),
            ("4", "Generate plan", "Confirm daily meals are returned."),
            ("5", "Check safety", "Verify dairy foods are excluded."),
            ("6", "Use results", "Export groceries and log food, water and weight."),
        ],
        "footer": "Live app: frontend-peach-nine-3tf89350e5.vercel.app   Health: backend-sigma-two-51.vercel.app/health",
    },
    {
        "title": "Practical AI systems need reliability and honesty",
        "subtitle": "Challenges and lessons learned",
        "items": [
            ("Allergy first", "Apply allergy filtering before scoring recommendations."),
            ("Fallbacks matter", "AI provider failure should not break the application."),
            ("Do not overclaim", "The photo feature is honest about low confidence."),
            ("Setup clarity", "Docker and environment variables must be documented."),
            ("End-to-end tests", "Meal planning touches profile, foods, allergies and groceries."),
            ("Explainability", "Safety-sensitive nutrition features need understandable decisions."),
        ],
    },
    {
        "type": "closing",
        "title": "NutriAI turns personal data into practical food decisions",
        "subtitle": "Conclusion",
        "items": ["Personalize", "Plan", "Shop", "Track", "Improve"],
        "footer": "One connected experience for nutrition planning, allergy awareness and AI-assisted guidance.",
    },
]


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = LINE


def add_textbox(slide, x, y, w, h, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title: str, subtitle: str, number: int) -> None:
    add_textbox(slide, 0.65, 0.28, 8.5, 0.3, subtitle.upper(), 10, GREEN_DARK, True)
    add_textbox(slide, 0.65, 0.62, 10.6, 0.7, title, 26, INK, True)
    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.8), Inches(0.32), Inches(0.38), Inches(0.38))
    set_fill(mark, GREEN)
    mark.line.color.rgb = GREEN
    add_textbox(slide, 11.88, 0.37, 0.22, 0.2, "N", 12, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.65, 7.04, 7.0, 0.25, "NutriAI - AI Nutrition and Meal Planner", 8, MUTED)
    add_textbox(slide, 11.65, 7.04, 0.45, 0.25, f"{number:02}", 8, GREEN, True, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=GREEN, soft=MINT) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, WHITE)
    card.line.color.rgb = LINE
    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(y + 0.25), Inches(0.48), Inches(0.48))
    set_fill(icon, soft)
    icon.line.color.rgb = soft
    add_textbox(slide, x + 0.34, y + 0.34, 0.25, 0.18, title[:1], 12, accent, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.85, y + 0.22, w - 1.05, 0.38, title, 14, INK, True)
    add_textbox(slide, x + 0.85, y + 0.68, w - 1.05, h - 0.8, body, 10.5, MUTED)


def add_grid_slide(prs, data, number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, data["title"], data["subtitle"], number)
    accents = [(GREEN, MINT), (BLUE, BLUE_SOFT), (AMBER, AMBER_SOFT), (PURPLE, PURPLE_SOFT), (CORAL, CORAL_SOFT), (GREEN_DARK, MINT)]
    items = data["items"]
    cols = 3 if len(items) > 4 else 2
    w = 3.55 if cols == 3 else 5.25
    x0 = 0.65 if cols == 3 else 0.9
    dx = 3.85 if cols == 3 else 5.7
    for i, item in enumerate(items):
        row, col = divmod(i, cols)
        title, body = item[-2], item[-1]
        accent, soft = accents[i % len(accents)]
        add_card(slide, x0 + col * dx, 1.55 + row * 1.95, w, 1.55, title, body, accent, soft)
    if "footer" in data:
        foot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(6.18), Inches(10.1), Inches(0.55))
        set_fill(foot, GREEN_DARK)
        foot.line.color.rgb = GREEN_DARK
        add_textbox(slide, 1.55, 6.34, 9.7, 0.2, data["footer"], 11, WHITE, True, PP_ALIGN.CENTER)


def add_cover(prs, data) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(7, 63, 49)
    add_textbox(slide, 0.75, 0.75, 0.6, 0.55, "N", 28, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.75, 1.75, 4.5, 0.3, "CLASS PRESENTATION", 11, RGBColor(116, 213, 172), True)
    add_textbox(slide, 0.75, 2.28, 7.2, 1.0, data["title"], 54, WHITE, True)
    add_textbox(slide, 0.78, 3.46, 6.5, 0.55, data["subtitle"], 25, RGBColor(207, 235, 221))
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(4.48), Inches(5.1), Inches(0.55))
    set_fill(tag, RGBColor(11, 81, 62))
    tag.line.color.rgb = RGBColor(35, 121, 91)
    add_textbox(slide, 1.0, 4.65, 4.65, 0.2, "Personalized meals - Safer choices - Smarter tracking", 11, WHITE, True, PP_ALIGN.CENTER)
    for i, value in enumerate(data["items"]):
        add_textbox(slide, 0.8 + i * 2.65, 6.2, 2.45, 0.25, value, 9.5, RGBColor(185, 216, 201), i == 0)


def add_closing(prs, data, number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(7, 63, 49)
    add_textbox(slide, 0.65, 0.32, 3.2, 0.25, data["subtitle"].upper(), 10, RGBColor(185, 216, 201), True)
    add_textbox(slide, 0.65, 0.72, 10.6, 0.85, data["title"], 26, WHITE, True)
    add_textbox(slide, 1.25, 2.05, 10.1, 0.55, data["footer"], 20, RGBColor(207, 235, 221), False, PP_ALIGN.CENTER)
    for i, item in enumerate(data["items"]):
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85 + i * 2.35), Inches(3.75), Inches(1.8), Inches(0.75))
        set_fill(pill, RGBColor(11, 81, 62))
        pill.line.color.rgb = RGBColor(35, 121, 91)
        add_textbox(slide, 0.95 + i * 2.35, 4.02, 1.6, 0.18, item, 12, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 4.8, 5.72, 3.0, 0.35, "Thank you", 24, RGBColor(105, 214, 168), True, PP_ALIGN.CENTER)
    add_textbox(slide, 11.65, 7.04, 0.45, 0.25, f"{number:02}", 8, RGBColor(105, 214, 168), True, PP_ALIGN.RIGHT)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(SLIDES, 1):
        if data.get("type") == "cover":
            add_cover(prs, data)
        elif data.get("type") == "closing":
            add_closing(prs, data, i)
        else:
            add_grid_slide(prs, data, i)

    prs.save(OUTPUT)


def main() -> None:
    build()
    print(f"Created {OUTPUT}")
    print(f"Slides: {len(SLIDES)}")
    print("Format: 16:9 widescreen")


if __name__ == "__main__":
    main()
